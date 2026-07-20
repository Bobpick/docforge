"""PPTX (PowerPoint) document handler — slides, text, shapes, tables, images."""

import re
import io
import zipfile
from pathlib import Path
from typing import Tuple, Dict, Any, List, Optional

from .utils import (
    format_table_md,
    format_table_json,
    detect_math_expression,
    clean_whitespace,
    generate_image_filename,
    is_duplicate_image,
)


class PPTXHandler:
    """Convert PPTX files to Markdown and structured JSON."""

    def __init__(self, extract_images: bool = True):
        self.extract_images = extract_images
        self._seen_image_hashes: set = set()

    def convert(
        self, file_path: Path
    ) -> Tuple[str, Dict[str, Any], List[Dict[str, Any]], Dict[str, Any]]:
        """Convert a PPTX file."""
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            raise ImportError(
                "python-pptx is required for PPTX conversion. "
                "Install it with: pip install python-pptx"
            )

        prs = Presentation(str(file_path))
        metadata = self._extract_metadata(prs, file_path)

        # Extract images from the ZIP
        images = []
        if self.extract_images:
            images = self._extract_images(file_path)

        sections: List[Dict[str, Any]] = []
        md_parts: List[str] = []

        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_md, slide_sections = self._process_slide(
                slide, slide_idx, len(images)
            )
            if slide_md.strip():
                md_parts.append(slide_md)
                sections.extend(slide_sections)

        markdown = clean_whitespace("\n\n---\n\n".join(md_parts))
        structured = self._build_structured(sections, metadata, images)
        return markdown, structured, images, metadata

    # ------------------------------------------------------------------
    # Metadata
    # ------------------------------------------------------------------
    def _extract_metadata(self, prs, file_path: Path) -> Dict[str, Any]:
        core = prs.core_properties
        return {
            "title": core.title or file_path.stem,
            "author": core.author or "",
            "subject": core.subject or "",
            "created": str(core.created) if core.created else "",
            "modified": str(core.modified) if core.modified else "",
            "slide_count": len(prs.slides),
        }

    # ------------------------------------------------------------------
    # Slide processing
    # ------------------------------------------------------------------
    def _process_slide(
        self, slide, slide_num: int, image_count: int
    ) -> Tuple[str, List[Dict[str, Any]]]:
        """Process a single slide."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        sections: List[Dict[str, Any]] = []
        md_lines: List[str] = []

        # Slide header
        md_lines.append(f"## Slide {slide_num}")
        sections.append(
            {"type": "heading", "level": 2, "content": f"Slide {slide_num}"}
        )

        for shape in slide.shapes:
            shape_type = shape.shape_type

            # --- Text boxes and shapes with text ---
            if shape.has_text_frame:
                text_md, text_sections = self._process_text_frame(
                    shape.text_frame, shape
                )
                if text_md.strip():
                    md_lines.append(text_md)
                    sections.extend(text_sections)

            # --- Tables ---
            if shape.has_table:
                table_md, table_sec = self._process_table(shape.table)
                if table_md.strip():
                    md_lines.append(table_md)
                    sections.append(table_sec)

            # --- Group shapes ---
            if shape_type == MSO_SHAPE_TYPE.GROUP:
                group_md, group_sections = self._process_group(shape, slide_num)
                if group_md.strip():
                    md_lines.append(group_md)
                    sections.extend(group_sections)

        return "\n\n".join(md_lines), sections

    def _process_text_frame(self, text_frame, shape) -> Tuple[str, List[Dict[str, Any]]]:
        """Process a text frame (title, body, text box)."""
        sections: List[Dict[str, Any]] = []
        md_lines: List[str] = []

        # Check if this is a title shape
        is_title = False
        try:
            ph = shape.placeholder_format
            if ph is not None:
                ph_type = ph.type
                # Title = 1, Center Title = 3
                if ph_type in (1, 3):
                    is_title = True
        except (ValueError, KeyError, AttributeError, TypeError):
            pass

        for para in text_frame.paragraphs:
            if not para.text.strip():
                continue

            text = self._format_paragraph_runs(para)
            level = para.level  # Indentation level

            if is_title and level == 0:
                md_lines.append(f"### {text}")
                sections.append({"type": "heading", "level": 3, "content": text})
            elif level > 0:
                indent = "  " * (level - 1)
                md_lines.append(f"{indent}- {text}")
                sections.append(
                    {"type": "list_item", "content": text, "level": level}
                )
            else:
                # Check for math
                if detect_math_expression(text):
                    md_lines.append(f"${text}$")
                    sections.append({"type": "math", "content": text})
                else:
                    md_lines.append(text)
                    sections.append({"type": "paragraph", "content": text})

        return "\n\n".join(md_lines), sections

    def _format_paragraph_runs(self, para) -> str:
        """Format paragraph runs with inline markdown."""
        parts = []
        for run in para.runs:
            text = run.text
            if not text:
                continue

            is_bold = run.font.bold
            is_italic = run.font.italic
            font_name = run.font.name or ""
            is_mono = any(
                m in (font_name or "").lower()
                for m in ["consol", "courier", "mono", "code"]
            )

            if is_mono:
                parts.append(f"`{text}`")
            elif is_bold and is_italic:
                parts.append(f"***{text}***")
            elif is_bold:
                parts.append(f"**{text}**")
            elif is_italic:
                parts.append(f"*{text}*")
            else:
                parts.append(text)

        return "".join(parts)

    def _process_table(self, table) -> Tuple[str, Dict[str, Any]]:
        """Process a PowerPoint table."""
        rows_data = []
        for row in table.rows:
            cells = [
                cell.text.strip().replace("\n", " ") for cell in row.cells
            ]
            rows_data.append(cells)

        if not rows_data:
            return "", {"type": "table", "headers": [], "rows": []}

        headers = rows_data[0]
        rows = rows_data[1:] if len(rows_data) > 1 else []

        md_text = format_table_md(headers, rows)
        return md_text, format_table_json(headers, rows)

    def _process_group(self, group, slide_num: int) -> Tuple[str, List[Dict[str, Any]]]:
        """Process a group shape recursively."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        sections: List[Dict[str, Any]] = []
        md_lines: List[str] = []

        for shape in group.shapes:
            if shape.has_text_frame:
                text_md, text_sections = self._process_text_frame(
                    shape.text_frame, shape
                )
                if text_md.strip():
                    md_lines.append(text_md)
                    sections.extend(text_sections)

            if shape.has_table:
                table_md, table_sec = self._process_table(shape.table)
                if table_md.strip():
                    md_lines.append(table_md)
                    sections.append(table_sec)

        return "\n\n".join(md_lines), sections

    # ------------------------------------------------------------------
    # Image extraction
    # ------------------------------------------------------------------
    def _extract_images(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract embedded images from a PPTX file (ZIP archive)."""
        images = []

        try:
            with zipfile.ZipFile(str(file_path), "r") as zf:
                media_files = [
                    f
                    for f in zf.namelist()
                    if f.startswith("ppt/media/") and not f.endswith("/")
                ]

                for idx, media_file in enumerate(media_files, 1):
                    img_data = zf.read(media_file)

                    if len(img_data) < 500:
                        continue
                    if is_duplicate_image(img_data, self._seen_image_hashes):
                        continue

                    ext = Path(media_file).suffix.lstrip(".") or "png"
                    if ext in ("emf", "wmf"):
                        ext = "png"

                    filename = f"image_{idx:03d}.{ext}"
                    images.append(
                        {
                            "filename": filename,
                            "data": img_data,
                            "extension": ext,
                            "page": 0,
                            "source": media_file,
                        }
                    )

        except Exception:
            pass

        return images

    # ------------------------------------------------------------------
    # Structured JSON
    # ------------------------------------------------------------------
    def _build_structured(
        self,
        sections: List[Dict[str, Any]],
        metadata: Dict[str, Any],
        images: List[Dict[str, Any]],
    ) -> Dict[str, Any]:
        """Build the structured JSON representation."""
        for img in images:
            sections.append(
                {
                    "type": "image",
                    "path": f"images/{img['filename']}",
                    "filename": img["filename"],
                }
            )

        return {
            "title": metadata.get("title", ""),
            "metadata": metadata,
            "sections": sections,
            "images": [
                {"path": f"images/{img['filename']}", "filename": img["filename"]}
                for img in images
            ],
            "stats": {
                "total_sections": len(sections),
                "headings": sum(1 for s in sections if s.get("type") == "heading"),
                "paragraphs": sum(1 for s in sections if s.get("type") == "paragraph"),
                "tables": sum(1 for s in sections if s.get("type") == "table"),
                "code_blocks": sum(1 for s in sections if s.get("type") == "code"),
                "math_blocks": sum(1 for s in sections if s.get("type") == "math"),
                "images": len(images),
            },
        }
